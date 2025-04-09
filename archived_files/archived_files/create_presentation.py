import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Set the style for visualizations
plt.style.use('seaborn-v0_8-whitegrid')
sns.set_palette("viridis")
plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['savefig.dpi'] = 300

# Create output directories
output_dir = 'presentation_materials'
temp_dir = os.path.join(output_dir, 'temp_images')
if not os.path.exists(output_dir):
    os.makedirs(output_dir)
if not os.path.exists(temp_dir):
    os.makedirs(temp_dir)

def load_data():
    """Load the DEA results"""
    dea_results = pd.read_csv('dea_results.csv')
    return dea_results

def generate_temp_images(data):
    """Generate temporary images for the presentation"""
    # Efficiency distribution
    plt.figure(figsize=(10, 6))
    sns.histplot(data=data, x='TE_VRS', bins=20, color='#3498db')
    plt.title('Distribution of Technical Efficiency (VRS)', fontsize=16, fontweight='bold')
    plt.xlabel('Technical Efficiency (VRS)', fontsize=12)
    plt.ylabel('Number of Hospitals', fontsize=12)
    plt.tight_layout()
    plt.savefig(os.path.join(temp_dir, 'efficiency_distribution.png'))
    plt.close()
    
    # Top hospitals
    top_hospitals = data.sort_values('TE_VRS', ascending=False).drop_duplicates('Hospital_Name').head(10)
    plt.figure(figsize=(12, 6))
    ax = sns.barplot(x='TE_VRS', y='Hospital_Name', data=top_hospitals, palette='viridis')
    plt.title('Top 10 Most Efficient Hospitals', fontsize=16, fontweight='bold')
    plt.xlabel('Technical Efficiency (VRS)', fontsize=12)
    plt.ylabel('Hospital', fontsize=12)
    # Add values to bars
    for i, v in enumerate(top_hospitals['TE_VRS']):
        ax.text(v + 0.1, i, f"{v:.2f}", va='center', fontweight='bold')
    plt.tight_layout()
    plt.savefig(os.path.join(temp_dir, 'top_hospitals.png'))
    plt.close()
    
    # Efficiency by categories
    fig, axs = plt.subplots(1, 3, figsize=(16, 5))
    sns.boxplot(x='Size_Category', y='TE_VRS', data=data, ax=axs[0], palette='viridis')
    axs[0].set_title('Efficiency by Size', fontsize=14, fontweight='bold')
    axs[0].set_xlabel('Hospital Size', fontsize=12)
    axs[0].set_ylabel('Technical Efficiency (VRS)', fontsize=12)
    
    sns.boxplot(x='Occupancy_Category', y='TE_VRS', data=data, ax=axs[1], palette='viridis')
    axs[1].set_title('Efficiency by Occupancy', fontsize=14, fontweight='bold')
    axs[1].set_xlabel('Occupancy Level', fontsize=12)
    axs[1].set_ylabel('Technical Efficiency (VRS)', fontsize=12)
    
    scale_palette = {'CRS': '#2ecc71', 'DRS': '#e74c3c', 'IRS': '#3498db'}
    sns.boxplot(x='Scale_Efficiency_Category', y='TE_VRS', data=data, ax=axs[2], palette=scale_palette)
    axs[2].set_title('Efficiency by Scale Category', fontsize=14, fontweight='bold')
    axs[2].set_xlabel('Scale Efficiency Category', fontsize=12)
    axs[2].set_ylabel('Technical Efficiency (VRS)', fontsize=12)
    
    plt.tight_layout()
    plt.savefig(os.path.join(temp_dir, 'efficiency_by_categories.png'))
    plt.close()
    
    # Scale efficiency pie chart
    scale_counts = data['Scale_Efficiency_Category'].value_counts()
    plt.figure(figsize=(8, 8))
    colors = ['#2ecc71', '#e74c3c', '#3498db']
    plt.pie(scale_counts, labels=scale_counts.index, autopct='%1.1f%%', 
            colors=colors, startangle=90, shadow=True)
    plt.title('Scale Efficiency Categories', fontsize=16, fontweight='bold')
    plt.tight_layout()
    plt.savefig(os.path.join(temp_dir, 'scale_efficiency_pie.png'))
    plt.close()
    
    # Correlation heatmap
    cols_to_include = ['TE_VRS', 'TE_CRS', 'Scale_Efficiency', 'Total_Surveys', 
                     'Patient_Satisfaction', 'Recommendation_Score', 
                     'Licensed_Beds', 'Staffed_Beds', 'Bed_Occupancy_Rate']
    correlation_data = data[cols_to_include].drop_duplicates()
    correlation_matrix = correlation_data.corr()
    
    plt.figure(figsize=(12, 10))
    sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f", 
               square=True, linewidths=.5, cbar_kws={"shrink": .8})
    plt.title('Correlation Matrix of Key Performance Indicators', fontsize=16, fontweight='bold')
    plt.tight_layout()
    plt.savefig(os.path.join(temp_dir, 'correlation_heatmap.png'))
    plt.close()
    
    # Efficiency vs quality scatter plot
    plt.figure(figsize=(12, 8))
    scatter = sns.scatterplot(data=data, x='Patient_Satisfaction', y='TE_VRS', 
                            hue='Size_Category', size='Bed_Occupancy_Rate',
                            sizes=(50, 250), alpha=0.7)
    plt.title('Efficiency vs Patient Satisfaction', fontsize=16, fontweight='bold')
    plt.xlabel('Patient Satisfaction', fontsize=12)
    plt.ylabel('Technical Efficiency (VRS)', fontsize=12)
    
    # Add trendline
    x = data['Patient_Satisfaction']
    y = data['TE_VRS']
    z = np.polyfit(x, y, 1)
    p = np.poly1d(z)
    plt.plot(x, p(x), "r--", alpha=0.7)
    
    # Add correlation coefficient
    corr = x.corr(y)
    plt.annotate(f"Correlation: {corr:.2f}", xy=(0.05, 0.95), xycoords='axes fraction',
                fontsize=12, bbox=dict(boxstyle="round,pad=0.3", fc="white", alpha=0.8))
    
    plt.legend(title='Hospital Size', bbox_to_anchor=(1.05, 1), loc='upper left')
    plt.tight_layout()
    plt.savefig(os.path.join(temp_dir, 'efficiency_vs_satisfaction.png'))
    plt.close()

def create_presentation(data):
    """Create a PowerPoint presentation with DEA results"""
    prs = Presentation()
    
    # Set slide dimensions and layouts
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Hospital Efficiency Analysis"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)
    
    subtitle.text = "Using Data Envelopment Analysis (DEA)"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    
    # Introduction slide
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introduction to Hospital Efficiency Analysis"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    content_text = [
        "This presentation summarizes the results of a comprehensive efficiency analysis of hospitals using Data Envelopment Analysis (DEA).",
        "",
        "Key Points:",
        "• Analyzed " + str(len(data.drop_duplicates('Hospital_Name'))) + " hospitals using DEA methodology",
        "• Technical Efficiency was calculated using both Variable Returns to Scale (VRS) and Constant Returns to Scale (CRS) models",
        "• Scale Efficiency was determined to identify optimal operating scales",
        "• Hospitals were categorized by size, occupancy rate, and quality measures",
        "• The relationship between efficiency and quality metrics was examined"
    ]
    
    content.text = "\n".join(content_text)
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        if i > 1:  # Apply bullet points starting from the third paragraph
            paragraph.level = 1 if i > 2 else 0
        paragraph.font.size = Pt(24) if i < 2 else Pt(20)
    
    # Key metrics slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Efficiency Metrics"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    avg_te_vrs = data['TE_VRS'].mean()
    avg_te_crs = data['TE_CRS'].mean()
    avg_scale_eff = data['Scale_Efficiency'].mean()
    
    metrics_text = [
        "Overview of Hospital Efficiency Analysis Results:",
        "",
        f"• Average Technical Efficiency (VRS): {avg_te_vrs:.2f}",
        f"• Average Technical Efficiency (CRS): {avg_te_crs:.2f}",
        f"• Average Scale Efficiency: {avg_scale_eff:.2f}",
        "",
        "Scale Efficiency Categories:",
        f"• Constant Returns to Scale (CRS): {data['Scale_Efficiency_Category'].value_counts().get('CRS', 0)} hospitals",
        f"• Increasing Returns to Scale (IRS): {data['Scale_Efficiency_Category'].value_counts().get('IRS', 0)} hospitals",
        f"• Decreasing Returns to Scale (DRS): {data['Scale_Efficiency_Category'].value_counts().get('DRS', 0)} hospitals"
    ]
    
    content.text = "\n".join(metrics_text)
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(24) if i < 2 else Pt(20)
        if i > 1:  # Apply bullet points starting from the third paragraph
            paragraph.level = 0
    
    # Efficiency distribution slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "Distribution of Technical Efficiency"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    img_path = os.path.join(temp_dir, 'efficiency_distribution.png')
    slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), width=Inches(10))
    
    # Top hospitals slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "Top 10 Most Efficient Hospitals"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    img_path = os.path.join(temp_dir, 'top_hospitals.png')
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12))
    
    # Efficiency by categories slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "Efficiency Analysis by Hospital Categories"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    img_path = os.path.join(temp_dir, 'efficiency_by_categories.png')
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12))
    
    # Scale efficiency pie chart slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "Scale Efficiency Categories"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    img_path = os.path.join(temp_dir, 'scale_efficiency_pie.png')
    slide.shapes.add_picture(img_path, Inches(2.5), Inches(1.5), width=Inches(8))
    
    # Correlation matrix slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "Correlation Matrix of Key Performance Indicators"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    img_path = os.path.join(temp_dir, 'correlation_heatmap.png')
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(11))
    
    # Efficiency vs quality slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "Relationship Between Efficiency and Patient Satisfaction"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    img_path = os.path.join(temp_dir, 'efficiency_vs_satisfaction.png')
    slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.5), width=Inches(11.5))
    
    # Key findings slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Findings & Recommendations"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    findings_text = [
        "Key Findings:",
        "• Small hospitals tend to show higher efficiency variance, with both the most and least efficient performers",
        "• Hospitals with high occupancy rates demonstrate better technical efficiency scores",
        f"• {data['Scale_Efficiency_Category'].value_counts().get('CRS', 0)/len(data)*100:.0f}% of hospitals operate at constant returns to scale (CRS), indicating optimal scale operations",
        "• Patient satisfaction correlates positively with technical efficiency",
        "",
        "Recommendations:",
        "• Focus improvement efforts on hospitals with decreasing returns to scale (DRS)",
        "• Implement best practices from top-performing small hospitals",
        "• Address operational inefficiencies in facilities with low bed occupancy rates",
        "• Develop standardized DEA-based performance monitoring systems",
        "• Conduct follow-up analysis to identify specific efficiency drivers"
    ]
    
    content.text = "\n".join(findings_text)
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(24) if i == 0 or i == 6 else Pt(20)
        if i > 0 and i != 6:  # Apply bullet points for all paragraphs except headers
            paragraph.level = 0
    
    # Conclusion slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    conclusion_text = [
        "Data Envelopment Analysis provides valuable insights into hospital efficiency:",
        "",
        "• DEA effectively evaluates hospital performance across multiple dimensions",
        "• The analysis reveals significant variation in efficiency across hospitals",
        "• Hospital size, occupancy, and scale efficiency are key factors affecting performance",
        "• Efficiency and quality metrics show positive correlation, suggesting that patient-centered care and operational efficiency can be achieved simultaneously",
        "",
        "Next Steps:",
        "• Expand analysis to include additional variables like financial performance",
        "• Develop benchmarking initiatives based on top-performing hospitals",
        "• Implement targeted efficiency improvement programs"
    ]
    
    content.text = "\n".join(conclusion_text)
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(24) if i < 2 or i == 7 else Pt(20)
        if i > 1 and i != 7:  # Apply bullet points appropriately
            paragraph.level = 0
    
    # Save the presentation
    ppt_path = os.path.join(output_dir, 'Hospital_Efficiency_Analysis.pptx')
    prs.save(ppt_path)
    print(f"PowerPoint presentation created at '{ppt_path}'")

def main():
    print("Generating presentation materials...")
    data = load_data()
    
    print("Generating temporary images...")
    generate_temp_images(data)
    
    print("Creating PowerPoint presentation...")
    create_presentation(data)
    
    print("Presentation creation completed.")

if __name__ == "__main__":
    main() 